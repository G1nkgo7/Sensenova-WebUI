"""Attachment storage and text extraction for Studio deck requests."""
import csv
import io
import json
import os
import re
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path

from fastapi import UploadFile

from . import engine

MAX_FILES = 8
MAX_FILE_BYTES = 20 * 1024 * 1024
MAX_TEXT_PER_FILE = 12000
MAX_TOTAL_TEXT = 50000
MAX_PDF_IMAGE_PAGES = 20

TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".jsonl", ".log",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".html", ".css", ".xml", ".yaml", ".yml",
}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff"}
MODE_WEB_PARSE = "web_parse"
MODE_PIPELINE_AGENT = "pipeline_agent"
MODES = {MODE_WEB_PARSE, MODE_PIPELINE_AGENT}


def normalize_mode(value: str | None = None) -> str:
    raw = (value or os.environ.get("STUDIO_ATTACHMENT_MODE") or MODE_WEB_PARSE).strip().lower()
    raw = raw.replace("-", "_")
    aliases = {
        "": MODE_WEB_PARSE,
        "web": MODE_WEB_PARSE,
        "parse": MODE_WEB_PARSE,
        "web_parse": MODE_WEB_PARSE,
        "pipeline": MODE_PIPELINE_AGENT,
        "raw": MODE_PIPELINE_AGENT,
        "agent": MODE_PIPELINE_AGENT,
        "pipeline_agent": MODE_PIPELINE_AGENT,
    }
    return aliases.get(raw, MODE_WEB_PARSE)


def safe_name(name: str) -> str:
    base = os.path.basename(name or "").replace("\\", "_").replace("/", "_").strip()
    base = re.sub(r"[^A-Za-z0-9._()\- \u4e00-\u9fff]+", "_", base)
    return base[:120] or "file"


def _unique_path(directory: Path, filename: str) -> Path:
    base = safe_name(filename)
    stem = Path(base).stem or "file"
    suffix = Path(base).suffix
    out = directory / base
    i = 2
    while out.exists():
        out = directory / f"{stem}_{i}{suffix}"
        i += 1
    return out


def _clean_text(text: str) -> str:
    text = text.replace("\x00", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _clip(text: str, limit: int = MAX_TEXT_PER_FILE) -> tuple[str, bool]:
    text = _clean_text(text)
    if len(text) <= limit:
        return text, False
    return text[:limit].rsplit("\n", 1)[0].strip() or text[:limit], True


def _decode_text(data: bytes) -> str:
    for enc in ("utf-8-sig", "utf-8", "gb18030", "big5", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_text_file(path: Path, data: bytes, ext: str) -> str:
    text = _decode_text(data)
    if ext in {".csv", ".tsv"}:
        dialect = csv.excel_tab if ext == ".tsv" else csv.excel
        rows = []
        try:
            for row in csv.reader(io.StringIO(text), dialect=dialect):
                rows.append(" | ".join(cell.strip() for cell in row))
        except Exception:
            return text
        return "\n".join(rows)
    if ext in {".json", ".jsonl"}:
        try:
            if ext == ".json":
                return json.dumps(json.loads(text), ensure_ascii=False, indent=2)
        except Exception:
            pass
    return text


def _asset_stem(path: Path) -> str:
    return re.sub(r"[^A-Za-z0-9._-]+", "_", path.stem)[:80] or "attachment"


def _image_asset(outdir: Path, stem: str, kind: str, idx: int, data: bytes,
                 suffix: str, page: int | None = None) -> dict | None:
    suffix = (suffix or ".png").lower()
    if suffix not in IMAGE_EXTS:
        return None
    rel = f"attachments/{stem}_{kind}_{idx:03d}{suffix}"
    out_path = outdir / Path(rel).name
    out_path.write_bytes(data)
    item = {
        "source_path": str(out_path),
        "workspace_rel": rel,
        "kind": kind,
    }
    if page is not None:
        item["page"] = page
    return item


def _extract_zip_media(data: bytes, prefix: str, outdir: Path, stem: str, kind: str) -> list[dict]:
    images = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            names = sorted(
                n for n in zf.namelist()
                if n.startswith(prefix) and Path(n).suffix.lower() in IMAGE_EXTS
            )
            for idx, name in enumerate(names, 1):
                item = _image_asset(outdir, stem, kind, idx, zf.read(name), Path(name).suffix)
                if item:
                    images.append(item)
    except Exception:
        pass
    return images


def _extract_docx(data: bytes, outdir: Path, stem: str) -> tuple[str, list[dict], list[str]]:
    text = _extract_docx_text(data)
    images = _extract_zip_media(data, "word/media/", outdir, stem, "docx_image")
    return text, images, []


def _extract_docx_text(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx_text(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            txt = getattr(shape, "text", "")
            if txt and txt.strip():
                texts.append(txt.strip())
        if texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


def _office_binary() -> str | None:
    return shutil.which("libreoffice") or shutil.which("soffice")


def _render_pdf_doc(doc, outdir: Path, stem: str, kind: str, limit: int | None = None) -> list[dict]:
    import fitz

    images = []
    n_pages = len(doc)
    if limit is not None:
        n_pages = min(n_pages, limit)
    for idx in range(n_pages):
        page = doc.load_page(idx)
        pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
        rel = f"attachments/{stem}_{kind}_{idx + 1:03d}.png"
        out_path = outdir / Path(rel).name
        pix.save(str(out_path))
        images.append({
            "source_path": str(out_path),
            "workspace_rel": rel,
            "page": idx + 1,
            "kind": kind,
        })
    return images


def _render_office_pages(path: Path, outdir: Path, stem: str, kind: str) -> tuple[list[dict], str | None]:
    office = _office_binary()
    if not office:
        return [], "LibreOffice/soffice 未安装,未生成页面截图。"
    try:
        with tempfile.TemporaryDirectory() as tmp:
            r = subprocess.run(
                [office, "--headless", "--convert-to", "pdf", "--outdir", tmp, str(path)],
                stdin=subprocess.DEVNULL,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                timeout=90,
                check=False,
            )
            pdfs = sorted(Path(tmp).glob("*.pdf"))
            if r.returncode != 0 or not pdfs:
                msg = (r.stderr or r.stdout or "未生成 PDF").strip()
                return [], f"LibreOffice 转 PDF 失败:{msg[:300]}"
            import fitz

            doc = fitz.open(str(pdfs[0]))
            return _render_pdf_doc(doc, outdir, stem, kind), None
    except Exception as e:
        return [], f"LibreOffice 渲染失败:{type(e).__name__}: {e}"


def _extract_pptx(path: Path, data: bytes, outdir: Path, stem: str) -> tuple[str, list[dict], list[str]]:
    text = _extract_pptx_text(data)
    images, notes = [], []
    rendered, note = _render_office_pages(path, outdir, stem, "pptx_slide")
    images.extend(rendered)
    if note:
        notes.append(note)
    images.extend(_extract_zip_media(data, "ppt/media/", outdir, stem, "pptx_image"))
    return text, images, notes


def _extract_xlsx_text(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        rows = []
        for r_i, row in enumerate(ws.iter_rows(values_only=True), 1):
            values = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if values:
                rows.append(" | ".join(values))
            if r_i >= 300:
                rows.append("[... sheet truncated after 300 rows ...]")
                break
        if rows:
            parts.append(f"[Sheet {ws.title}]\n" + "\n".join(rows))
    return "\n\n".join(parts)


def _extract_xlsx_images(data: bytes, outdir: Path, stem: str) -> tuple[list[dict], bool]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data))
    images = []
    has_charts = False
    idx = 0
    for ws in wb.worksheets:
        has_charts = has_charts or bool(getattr(ws, "_charts", []))
        for img in getattr(ws, "_images", []):
            idx += 1
            try:
                raw = img._data()
            except Exception:
                continue
            suffix = Path(getattr(img, "path", "") or "").suffix or ".png"
            item = _image_asset(outdir, stem, "xlsx_image", idx, raw, suffix)
            if item:
                item["sheet"] = ws.title
                images.append(item)
    return images, has_charts


def _extract_xlsx(path: Path, data: bytes, outdir: Path, stem: str) -> tuple[str, list[dict], list[str]]:
    text = _extract_xlsx_text(data)
    images, has_charts = _extract_xlsx_images(data, outdir, stem)
    notes = []
    rendered, note = _render_office_pages(path, outdir, stem, "xlsx_sheet")
    images.extend(rendered)
    if note and has_charts:
        notes.append("XLSX 中检测到图表;" + note)
    return text, images, notes


def _extract_pdf_text(data: bytes) -> list[str]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"[Page {i}]\n{text.strip()}")
    return pages


def _render_pdf_pages(data: bytes, outdir: Path, stem: str) -> list[dict]:
    import fitz

    doc = fitz.open(stream=data, filetype="pdf")
    return _render_pdf_doc(doc, outdir, stem, "pdf_page", MAX_PDF_IMAGE_PAGES)


def _image_info(path: Path) -> str:
    from PIL import Image

    with Image.open(path) as im:
        return f"{im.format or path.suffix.upper().lstrip('.')} image, {im.width}x{im.height}px"


def _parse_saved_file(path: Path, data: bytes, original_name: str) -> dict:
    ext = path.suffix.lower()
    record = {
        "name": original_name,
        "stored_name": path.name,
        "path": str(path),
        "size": len(data),
        "type": ext.lstrip(".") or "unknown",
        "text": "",
        "truncated": False,
        "images": [],
        "notes": [],
    }

    try:
        if ext in TEXT_EXTS:
            record["text"], record["truncated"] = _clip(_extract_text_file(path, data, ext))
        elif ext == ".docx":
            text, images, notes = _extract_docx(data, path.parent, _asset_stem(path))
            record["text"], record["truncated"] = _clip(text)
            record["images"] = images
            record["notes"].extend(notes)
        elif ext == ".pptx":
            text, images, notes = _extract_pptx(path, data, path.parent, _asset_stem(path))
            record["text"], record["truncated"] = _clip(text)
            record["images"] = images
            record["notes"].extend(notes)
        elif ext in {".xlsx", ".xlsm"}:
            text, images, notes = _extract_xlsx(path, data, path.parent, _asset_stem(path))
            record["text"], record["truncated"] = _clip(text)
            record["images"] = images
            record["notes"].extend(notes)
        elif ext == ".pdf":
            pages = _extract_pdf_text(data)
            text = "\n\n".join(pages)
            if _clean_text(text):
                record["text"], record["truncated"] = _clip(text)
            else:
                stem = re.sub(r"[^A-Za-z0-9._-]+", "_", path.stem)[:80] or "pdf"
                record["images"] = _render_pdf_pages(data, path.parent, stem)
                record["notes"].append("PDF 未提取到足够可复制文本,已按页转换为图片。")
                if record["images"]:
                    record["text"] = "扫描版 PDF,已将页面图片直接附给模型。"
        elif ext in IMAGE_EXTS:
            record["text"] = _image_info(path)
            record["images"] = [{
                "source_path": str(path),
                "workspace_rel": f"attachments/{path.name}",
                "kind": "image",
            }]
        else:
            record["notes"].append("暂不支持解析该文件类型,仅保留文件元信息。")
    except Exception as e:
        record["notes"].append(f"解析失败:{type(e).__name__}: {e}")
    return record


def _raw_saved_file(path: Path, data: bytes, original_name: str) -> dict:
    ext = path.suffix.lower()
    return {
        "name": original_name,
        "stored_name": path.name,
        "path": str(path),
        "size": len(data),
        "type": ext.lstrip(".") or "unknown",
        "text": "",
        "truncated": False,
        "images": [],
        "notes": ["原始附件已保存,将交给管线侧 attachment_reader 处理。"],
    }


def _record_block(idx: int, record: dict) -> str:
    lines = [
        f"{idx}. 文件名: {record['name']}",
        f"类型: {record.get('type') or 'unknown'}; 大小: {record.get('size', 0)} bytes",
    ]
    for note in record.get("notes") or []:
        lines.append(f"说明: {note}")
    images = record.get("images") or []
    if images:
        lines.append("图片附件/文档视觉页:")
        lines.extend(f"- {img['workspace_rel']}" for img in images)
        lines.append("这些图片会以二进制 image block 直接附在初始模型消息中,请结合图片内容理解附件。")
    text = record.get("text") or ""
    if text:
        label = "提取文本"
        if record.get("truncated"):
            label += "(已截断)"
        lines.append(f"{label}:\n{text}")
    return "\n".join(lines)


def _raw_rel(record: dict) -> str:
    stored = safe_name(record.get("stored_name") or record.get("name") or "file")
    return f"attachments/raw/{stored}"


def attachment_query_block(records: list[dict]) -> str:
    if not records:
        return ""
    blocks, total = [], 0
    for i, record in enumerate(records, 1):
        block = _record_block(i, record)
        remaining = MAX_TOTAL_TEXT - total
        if remaining <= 0:
            blocks.append("[... 附件内容已达到总长度上限,后续附件仅保留在 seed_json 元信息中 ...]")
            break
        if len(block) > remaining:
            block = block[:remaining].rsplit("\n", 1)[0].strip() + "\n[... 附件内容达到总长度上限,已截断 ...]"
        blocks.append(block)
        total += len(block)
    return "\n\n【附件资料】\n" + "\n\n".join(blocks)


def pipeline_agent_query_block(records: list[dict]) -> str:
    if not records:
        return ""
    lines = [
        "",
        "",
        "【附件处理模式:pipeline_agent】",
        "本请求包含原始附件;web 端没有预先解析正文,也不要假设附件内容已经出现在 prompt 中。",
        "原始附件会在运行工作区 staged 到 `attachments/raw/`,清单为 `attachments/manifest.json`。",
        "请在正式规划 PPT 前,先委派一个 `attachment_reader` 子 agent,toolsets 使用 `file`, `terminal`, `vision`,",
        "让它读取 `attachments/manifest.json` 和原始附件,必要时用 python/LibreOffice/pypdf/python-pptx/openpyxl/PIL 等工具解析,",
        "对图片、扫描页、PPT 页面截图或文档内嵌图使用 `vision_analyze`,最后把结构化结论交回编排器。",
        "附件列表:",
    ]
    for i, record in enumerate(records, 1):
        lines.append(
            f"{i}. {record.get('name') or record.get('stored_name') or 'file'} "
            f"({record.get('type') or 'unknown'}, {record.get('size', 0)} bytes) -> `{_raw_rel(record)}`"
        )
    return "\n".join(lines)


def store_attachments(deck_id: int, files: list[UploadFile], parse: bool = True) -> list[dict]:
    out = []
    updir = engine.deck_uploads_dir(deck_id)
    updir.mkdir(parents=True, exist_ok=True)
    for f in files or []:
        if len(out) >= MAX_FILES:
            break
        if not getattr(f, "filename", ""):
            continue
        data = f.file.read()
        if not data:
            continue
        if len(data) > MAX_FILE_BYTES:
            out.append({
                "name": f.filename,
                "stored_name": "",
                "path": "",
                "size": len(data),
                "type": Path(f.filename).suffix.lower().lstrip(".") or "unknown",
                "text": "",
                "truncated": False,
                "images": [],
                "notes": [f"文件超过 {MAX_FILE_BYTES // 1024 // 1024}MB,已跳过解析和保存。"],
            })
            continue
        path = _unique_path(updir, f.filename)
        path.write_bytes(data)
        out.append(_parse_saved_file(path, data, f.filename) if parse else _raw_saved_file(path, data, f.filename))
    return out


def store_and_parse(deck_id: int, files: list[UploadFile]) -> list[dict]:
    return store_attachments(deck_id, files, parse=True)


def attachment_images(records: list[dict]) -> list[dict]:
    images = []
    for record in records or []:
        for img in record.get("images") or []:
            if img.get("source_path") and img.get("workspace_rel"):
                images.append({
                    "source_path": img["source_path"],
                    "workspace_rel": img["workspace_rel"],
                    "name": record.get("name", ""),
                    "kind": img.get("kind", "image"),
                    "page": img.get("page"),
                })
    return images


def raw_attachments(records: list[dict]) -> list[dict]:
    raw = []
    for record in records or []:
        path = record.get("path")
        if not path:
            continue
        raw.append({
            "source_path": path,
            "workspace_rel": _raw_rel(record),
            "name": record.get("name", ""),
            "stored_name": record.get("stored_name", ""),
            "type": record.get("type", "unknown"),
            "size": record.get("size", 0),
        })
    return raw

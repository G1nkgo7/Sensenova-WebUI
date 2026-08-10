"""成稿图 → 图片型 PPTX(每页一张全幅 render,16:9)。"""
import glob
import os

from pptx import Presentation
from pptx.util import Emu

W, H = 12192000, 6858000      # 13.333in x 7.5in(与 1600x900 成稿等比)


def renders_of(run_dir: str):
    return sorted(glob.glob(os.path.join(run_dir, "renders", "slide_*.png")))


def build_pptx(run_dir: str, out_path: str) -> int:
    pngs = renders_of(run_dir)
    if not pngs:
        raise ValueError("该 deck 没有成稿图")
    prs = Presentation()
    prs.slide_width = Emu(W)
    prs.slide_height = Emu(H)
    blank = prs.slide_layouts[6]
    for p in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(p, 0, 0, width=prs.slide_width, height=prs.slide_height)
    tmp = out_path + ".tmp"
    prs.save(tmp)
    os.replace(tmp, out_path)
    return len(pngs)
